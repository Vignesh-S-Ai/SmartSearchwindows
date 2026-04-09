# extractor/document_parser.py
# Phase 2: File Watcher & Crawler
# This module parses various document formats (PDF, DOCX, TXT, etc.)

import csv
import json
from pathlib import Path

from config.logger import get_logger
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

logger = get_logger(__name__)

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB limit


def _read_text_file(path: Path) -> str:
    """Read plain text file with encoding fallback."""
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]

    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    # Last resort: read as binary and decode with error ignoring
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="ignore")


def _parse_pdf(path: Path) -> str:
    """Extract text from PDF file using PyPDF2."""
    text_parts = []
    try:
        with open(path, "rb") as f:
            reader = PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                except Exception as e:
                    logger.warning(
                        f"Failed to extract text from PDF page {page_num + 1} in {path}: {e}"
                    )
                    continue

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to parse PDF {path}: {e}")
        return ""


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX file using python-docx."""
    try:
        doc = Document(path)
        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to parse DOCX {path}: {e}")
        return ""


def _parse_pptx(path: Path) -> str:
    """Extract text from PowerPoint file using python-pptx."""
    try:
        prs = Presentation(path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides):
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            except Exception as e:
                logger.warning(
                    f"Failed to extract text from slide {slide_num + 1} in {path}: {e}"
                )
                continue

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to parse PPTX {path}: {e}")
        return ""


def _parse_xlsx(path: Path) -> str:
    """Extract text from Excel file using openpyxl."""
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(str(cell.value))

        wb.close()
        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to parse XLSX {path}: {e}")
        return ""


def _parse_csv(path: Path) -> str:
    """Parse CSV file and return as text."""
    try:
        text_parts = []
        with open(path, "r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                text_parts.append(
                    " | ".join(str(cell) for cell in row if cell is not None)
                )
        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to parse CSV {path}: {e}")
        return ""


def _parse_json(path: Path) -> str:
    """Parse JSON file and return as formatted text."""
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2)
    except Exception as e:
        logger.warning(f"Failed to parse JSON {path}: {e}")
        return ""


def _parse_yaml(path: Path) -> str:
    """Parse YAML file and return as text (fallback to raw read)."""
    try:
        # Try to import yaml
        import yaml

        with open(path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)

        if data:
            return json.dumps(data, indent=2)
        return ""
    except ImportError:
        # If PyYAML not installed, fall back to raw text read
        return _read_text_file(path)
    except Exception as e:
        logger.warning(f"Failed to parse YAML {path}: {e}")
        return ""


def parse_document(path: Path) -> str:
    """
    Parse a document file and extract text content.

    Supports: PDF, DOCX, DOC, PPTX, XLSX, CSV, JSON, YAML, TXT, MD,
    and common code files (PY, JS, TS, HTML, CSS).

    Args:
        path: Path to the document file

    Returns:
        Extracted text content as string, or empty string if extraction fails
    """
    # Check if file exists
    if not path.exists():
        logger.warning(f"File does not exist: {path}")
        return ""

    # Check file size
    try:
        file_size = path.stat().st_size
        if file_size > MAX_FILE_SIZE:
            logger.info(
                f"Skipped file {path} (file too large: {file_size / 1024 / 1024:.1f}MB > 10MB)"
            )
            return ""
        if file_size == 0:
            logger.info(f"Skipped file {path} (empty file)")
            return ""
    except Exception as e:
        logger.warning(f"Could not get file size for {path}: {e}")

    # Get the file extension in lowercase
    ext = path.suffix.lower()

    # Map extensions to parser functions
    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".csv": _parse_csv,
        ".json": _parse_json,
        ".yaml": _parse_yaml,
        ".yml": _parse_yaml,
    }

    # Check if this is a text-based file
    text_extensions = {
        ".txt",
        ".md",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".html",
        ".htm",
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".xml",
        ".sql",
        ".sh",
        ".bat",
        ".ps1",
        ".go",
        ".rs",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".cs",
        ".rb",
        ".php",
        ".swift",
        ".kt",
        ".scala",
        ".r",
        ".m",
        ".log",
        ".ini",
        ".cfg",
        ".conf",
        ".toml",
        ".env",
        ".gitignore",
        ".dockerfile",
        ".makefile",
    }

    try:
        # Try to find a parser for this extension
        if ext in parsers:
            text = parsers[ext](path)
            if text and text.strip():
                logger.debug(f"Successfully parsed {path} ({ext})")
                return text

        # Handle text-based files
        if ext in text_extensions:
            text = _read_text_file(path)
            if text and text.strip():
                logger.debug(f"Successfully read text file {path} ({ext})")
                return text

        # Extension not supported
        logger.info(f"Skipped file {path} (unsupported type: {ext})")
        return ""

    except Exception as e:
        logger.error(f"Error parsing file {path}: {e}")
        return ""
